"""
PowerPoint presentation generator for Delagrave product tech sheets.

This module provides the core PowerPoint slide generation engine that produces
product tech sheet slides from MD reference data. It handles template loading,
layout selection, text placeholder population, and image insertion.

This is a pure library module (no I/O beyond file operations).
"""

from pathlib import Path
from typing import Dict, List, Any, Optional
import tempfile
import zipfile
import shutil
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

from gendoc.parsers.md_parser import find_product, find_product_pages

# Delagrave corporate blue for arrow bullets
DELAGRAVE_BLUE = RGBColor(0, 85, 164)


# Family to layout index mapping
FAMILY_LAYOUT_MAP = {
    'paillasse': 1,        # Fiche technique profil paillasse
    'sorbonne': 2,         # Fiche technique sorbonne
    'revetement': 3,       # Fiche technique revetement
    'meubles': 4,          # Fiche technique meuble
    'tables-en': 4,        # Same layout as meubles
    'equipement': 5,       # Fiche technique equipement
    'elec-sorb': 5,        # Same layout as equipement
    'complements': 5,      # Same layout as equipement
    'armoire-securite': 0, # Modern template only (2 pages per product)
    'enceinte-ventilee': 0, # Modern template only (2 pages per product)
}

# VBA shape index to python-pptx placeholder index mapping
VBA_TO_PLACEHOLDER = {
    'paillasse': {
        1: 13,  # TEXTE
        2: 0,   # TITRE
        3: 14,  # Revetements
        4: 15,  # REF
        5: 16,  # dim 1
        6: 17,  # dim 2
        7: 18,  # dim 3
        8: 19,  # dim 4
        9: 20,  # dim 5
        10: 21, # page
    },
    'sorbonne': {
        1: 13,  # TEXTE
        2: 0,   # TITRE
        3: 14,  # Revetements
        4: 15,  # REF
        5: 16,  # dim 1
        6: 17,  # dim 2
        7: 18,  # dim 3
        8: 19,  # dim 4
        9: 20,  # dim 5
        10: 23, # dim 6
        11: 28, # dim 7
        12: 29, # page
    },
    'revetement': {
        1: 13,  # TEXTE
        2: 0,   # TITRE
        3: 14,  # Mise en oeuvre / Applications
        4: 15,  # REF
        5: 16,  # Finition
        6: 17,  # Applications
        7: 18,  # page
    },
    'meubles': {
        1: 13,  # TEXTE
        2: 0,   # TITRE
        3: 15,  # REF
        4: 16,  # page
    },
    'equipement': {
        1: 0,   # TITRE
        2: 15,  # Reference
        3: 16,  # page
    },
    'tables-en': {
        1: 13,  # TEXTE
        2: 0,   # TITRE
        3: 15,  # REF
        4: 16,  # page
    },
    'elec-sorb': {
        1: 0,   # TITRE
        2: 15,  # Reference
        3: 16,  # page
    },
    'complements': {
        1: 0,   # TITRE
        2: 15,  # Reference
        3: 16,  # page
    },
}


def _strip_template_shapes(prs: Presentation) -> None:
    """Remove all shapes from slide master and layouts to prevent template artifacts."""
    shape_tags = {'sp', 'pic', 'grpSp', 'cxnSp'}
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

    for master in prs.slide_masters:
        spTree = master._element.find('.//p:cSld/p:spTree', nsmap)
        if spTree is not None:
            for child in list(spTree):
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag in shape_tags:
                    spTree.remove(child)

    for layout in prs.slide_layouts:
        spTree = layout._element.find('.//p:cSld/p:spTree', nsmap)
        if spTree is not None:
            for child in list(spTree):
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag in shape_tags:
                    spTree.remove(child)


def load_template(template_path: Path) -> Presentation:
    """
    Load PowerPoint template from .potm (macro-enabled) file.

    python-pptx cannot open .potm directly, so this function:
    1. Copies .potm to temp file with .pptx extension
    2. Opens as zipfile and modifies content types
    3. Removes VBA macros
    4. Returns valid Presentation object

    Args:
        template_path: Path to the .potm template file

    Returns:
        Presentation object ready for slide generation

    Raises:
        FileNotFoundError: If template_path does not exist
        ValueError: If template cannot be converted
    """
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    # Create temporary directory for conversion
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)

        # Copy template to temp location
        temp_potm = temp_path / "template.potm"
        shutil.copy2(template_path, temp_potm)

        # Create output pptx path
        temp_pptx = temp_path / "template.pptx"

        # Convert .potm to .pptx by modifying zip contents
        with zipfile.ZipFile(temp_potm, 'r') as zip_in:
            with zipfile.ZipFile(temp_pptx, 'w', zipfile.ZIP_DEFLATED) as zip_out:
                for item in zip_in.infolist():
                    # Skip VBA project file
                    if 'vbaProject.bin' in item.filename:
                        continue

                    # Read file content
                    data = zip_in.read(item.filename)

                    # Modify [Content_Types].xml to remove macro references
                    if item.filename == '[Content_Types].xml':
                        data_str = data.decode('utf-8')
                        # Remove macro content type
                        data_str = data_str.replace(
                            'application/vnd.ms-powerpoint.template.macroEnabled.main+xml',
                            'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
                        )
                        # Remove vbaProject content type entries
                        data_str = data_str.replace(
                            '<Default Extension="bin" ContentType="application/vnd.ms-office.vbaProject"/>',
                            ''
                        )
                        data = data_str.encode('utf-8')

                    # Write to output zip
                    zip_out.writestr(item, data)

        # Load the converted presentation
        prs = Presentation(str(temp_pptx))

        return prs


def _split_revetement_text(full_text: str) -> Dict[str, str]:
    """
    Split revetement full text into TEXTE, MISE_EN_OEUVRE, and FINITION sections.

    Revetement texts consistently follow this structure:
    - Block 1 (before first blank line): Product description
    - Block 2 (after first blank line): Manufacturing/installation method
    - Block 3+ (after second blank line): Finish, colors, options

    Args:
        full_text: Full product text with newline separators

    Returns:
        Dict with keys 'texte', 'mise_en_oeuvre', 'finition'
    """
    blocks = [b.strip() for b in full_text.split('\n\n') if b.strip()]

    result = {
        'texte': '',
        'mise_en_oeuvre': '',
        'finition': ''
    }

    if len(blocks) >= 3:
        result['texte'] = blocks[0]
        result['mise_en_oeuvre'] = blocks[1]
        result['finition'] = '\n\n'.join(blocks[2:])
    elif len(blocks) == 2:
        result['texte'] = blocks[0]
        result['mise_en_oeuvre'] = blocks[1]
    elif len(blocks) == 1:
        result['texte'] = blocks[0]

    return result


def _format_text_with_arrows(placeholder: Any, text_content: str) -> None:
    """
    Format text placeholder with blue arrow prefix on each line and justified alignment.

    Each non-empty line gets a separate paragraph with:
    - Blue "▸ " arrow prefix
    - Justified text alignment
    - Compact line spacing

    Args:
        placeholder: python-pptx placeholder shape
        text_content: Raw text with newline separators
    """
    lines = [l.strip() for l in text_content.split('\n') if l.strip()]
    if not lines:
        return

    tf = placeholder.text_frame
    tf.text = ""  # Clear existing content (creates one empty paragraph)

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Blue arrow prefix
        arrow = p.add_run()
        arrow.text = "\u25B8 "  # ▸ small right-pointing triangle
        arrow.font.color.rgb = DELAGRAVE_BLUE

        # Text content
        text_run = p.add_run()
        text_run.text = line

        # Paragraph formatting
        p.alignment = PP_ALIGN.JUSTIFY
        p.space_after = Pt(1)

    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _populate_slide(slide: Any, product: Dict[str, Any], family: str) -> None:
    """
    Populate slide text placeholders with product data.

    Args:
        slide: python-pptx Slide object
        product: Product dictionary from md_parser
        family: Family name (lowercase)
    """
    # Create mapping of placeholder idx to content
    placeholder_data = {}

    # Map standard fields
    placeholder_data[0] = product.get('titre', '')
    placeholder_data[13] = product.get('texte', '')
    placeholder_data[15] = product.get('ref', '')

    # Filter sentinel "Aucune" value (means no text in source data)
    if placeholder_data.get(13, '').strip().lower() == 'aucune':
        placeholder_data[13] = ''

    # For revetement, split texte into TEXTE, MISE_EN_OEUVRE, FINITION
    if family == 'revetement' and product.get('texte', ''):
        text_parts = _split_revetement_text(product['texte'])
        placeholder_data[13] = text_parts['texte']          # TEXTE
        placeholder_data[16] = text_parts['mise_en_oeuvre']  # MISE_EN_OEUVRE
        placeholder_data[17] = text_parts['finition']        # FINITION

    # Map dimensions to placeholders using VBA_TO_PLACEHOLDER
    vba_mapping = VBA_TO_PLACEHOLDER.get(family, {})

    for dimension in product.get('dimensions', []):
        shape_index_str = dimension.get('shape_index', '').strip()
        if not shape_index_str:
            continue

        try:
            vba_index = int(shape_index_str)
        except (ValueError, TypeError):
            continue

        # Get placeholder idx from VBA mapping
        placeholder_idx = vba_mapping.get(vba_index)
        if placeholder_idx is None:
            continue

        # Build text value with prefix if present
        prefix = dimension.get('prefix', '').strip()
        valeur = dimension.get('valeur', '').strip()

        if prefix:
            text_value = f"{prefix}{valeur}"
        else:
            text_value = valeur

        placeholder_data[placeholder_idx] = text_value

    # Collect all placeholder indices used by this family
    family_placeholder_indices = set(vba_mapping.values())
    family_placeholder_indices.update([0, 13, 15])  # standard fields: titre, texte, ref

    # Determine which placeholders get blue arrow formatting
    arrow_indices = {13}  # TEXTE always gets arrows
    if family == 'revetement':
        arrow_indices.update({16, 17})  # mise_en_oeuvre and finition too

    # Populate placeholders (and remove empty ones to hide default template prompt)
    placeholders_to_remove = []
    for placeholder in slide.placeholders:
        try:
            idx = placeholder.placeholder_format.idx
            if idx in placeholder_data and placeholder_data[idx]:
                if idx in arrow_indices:
                    _format_text_with_arrows(placeholder, placeholder_data[idx])
                else:
                    tf = placeholder.text_frame
                    tf.text = placeholder_data[idx]
                    # Enable auto-shrink so text fits within the placeholder bounds
                    tf.word_wrap = True
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            elif idx in family_placeholder_indices:
                # Mark for removal to hide "Cliquez pour ajouter du texte" prompt
                placeholders_to_remove.append(placeholder)
        except (AttributeError, KeyError):
            # Some shapes may not have placeholder_format or idx
            continue

    # Remove empty placeholders (done after iteration to avoid modifying collection)
    for placeholder in placeholders_to_remove:
        sp = placeholder._element
        sp.getparent().remove(sp)


def _insert_images(slide: Any, product: Dict[str, Any], project_root: Path, family: str = '') -> int:
    """
    Insert product images into slide at specified positions.

    Args:
        slide: python-pptx Slide object
        product: Product dictionary from md_parser
        project_root: Root directory for resolving relative image paths
        family: Family name for family-specific adjustments

    Returns:
        Number of images successfully inserted
    """
    images_inserted = 0

    for image_data in product.get('images', []):
        chemin = image_data.get('chemin', '').strip()
        if not chemin:
            continue

        # Skip .missing files
        if chemin.endswith('.missing'):
            continue

        # Build absolute path
        image_path = project_root / chemin

        # Check if image exists
        if not image_path.exists():
            continue

        # Get position data (values are in points - VBA PowerPoint standard unit)
        left = image_data.get('left', 0)
        top = image_data.get('top', 0)
        width = image_data.get('width', 0)
        height = image_data.get('height', 0)

        if left == 0 or top == 0 or width == 0:
            continue

        # Convert points to EMUs (1 point = 1/72 inch = 12700 EMUs)
        EMU_PER_POINT = 12700
        left_emu = Emu(int(left * EMU_PER_POINT))
        top_emu = Emu(int(top * EMU_PER_POINT))
        width_emu = Emu(int(width * EMU_PER_POINT))

        # Reduce paillasse images slightly to give more space to description
        if family == 'paillasse':
            original_width = width_emu
            width_emu = Emu(int(width_emu * 0.85))
            # Shift image right to compensate
            left_emu = Emu(int(left_emu + (original_width - width_emu) * 0.5))

        # Insert image, clamping to slide bounds if needed
        # Slide dimensions: A4 portrait = 7561263 x 10693400 EMU
        SLIDE_WIDTH = 7561263
        SLIDE_HEIGHT = 10693400

        try:
            if height > 0:
                height_emu = Emu(int(height * EMU_PER_POINT))
            else:
                # Auto-calculate height from image aspect ratio
                from PIL import Image as PILImage
                with PILImage.open(str(image_path)) as im:
                    img_w, img_h = im.size
                ratio = img_h / img_w if img_w > 0 else 1.0
                height_emu = Emu(int(width_emu * ratio))

            # Clamp: if image overflows right, reduce width (keep ratio)
            if left_emu + width_emu > SLIDE_WIDTH:
                width_emu = Emu(SLIDE_WIDTH - left_emu)
                height_emu = Emu(int(width_emu * ratio)) if height <= 0 else height_emu

            # Clamp: if image overflows bottom, reduce height (keep ratio)
            if top_emu + height_emu > SLIDE_HEIGHT:
                max_h = Emu(SLIDE_HEIGHT - top_emu)
                scale = max_h / height_emu
                height_emu = max_h
                width_emu = Emu(int(width_emu * scale))

            slide.shapes.add_picture(
                str(image_path),
                left_emu,
                top_emu,
                width_emu,
                height_emu
            )
            images_inserted += 1
        except Exception:
            # Image may be corrupted or invalid format
            continue

    return images_inserted


def _add_revetement_slides(
    prs: Presentation,
    revetement_codes: List[str],
    references_dir: Path,
    project_root: Path,
    logo_path: Optional[Path] = None
) -> List[str]:
    """
    Add revetement (coating) slides for specified coating codes.

    Args:
        prs: Presentation object to add slides to
        revetement_codes: List of coating codes to generate slides for
        references_dir: Path to Delagrave/references/
        project_root: Root directory for resolving image paths
        logo_path: Optional path to logo image

    Returns:
        List of coating codes that were successfully added
    """
    from gendoc.generators.modern_template import build_product_slide

    added = []

    for code in revetement_codes:
        # Look up coating product
        product = find_product(code, references_dir)
        if not product:
            continue

        build_product_slide(prs, product, 'revetement', project_root, logo_path)
        added.append(code)

    return added


def generate_presentation(
    product_codes: List[str],
    output_path: Path,
    references_dir: Path,
    template_path: Path,
    project_root: Path,
    mode: str = "FTI",
    revetement_codes: Optional[List[str]] = None,
    devis_info: Optional[Dict[str, str]] = None,
    custom_products: Optional[List[Dict]] = None
) -> Dict[str, Any]:
    """
    Generate a PowerPoint presentation with product tech sheets.

    This is the main orchestrator function that:
    1. Loads the PowerPoint template
    2. Groups products by family
    3. Assembles complete document with cover, TOC, chapter separators, and product slides
    4. Auto-generates coating slides for products with coatings
    5. Saves the final presentation

    Args:
        product_codes: List of product codes to generate slides for
        output_path: Path where .pptx should be saved
        references_dir: Path to Delagrave/references/
        template_path: Path to .potm template file
        project_root: Root directory for resolving relative paths
        mode: Generation mode (default "FTI")
        revetement_codes: Optional list of coating codes to add manually
        devis_info: Optional dict with 'numero_devis', 'date', 'client', 'titre_affaire'
        custom_products: Optional list of custom product dicts (for SP articles)

    Returns:
        Dictionary with generation results:
        {
            'slides_generated': int,
            'total_pages': int,
            'revetements_added': list of coating codes,
            'skipped': list of {'code': str, 'reason': str}
        }

    Raises:
        FileNotFoundError: If template or references not found
    """
    # Lazy import to avoid circular dependency at module load time
    from gendoc.generators.document_assembler import assemble_document, FAMILY_ORDER
    from gendoc.generators.modern_template import init_company_info
    from collections import OrderedDict

    # Load company info from _entreprise.md
    init_company_info(references_dir)

    # Load template and remove any pre-existing blank slides
    prs = load_template(template_path)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # Strip template master/layout shapes (modern template builds all visuals)
    _strip_template_shapes(prs)

    skipped = []
    coating_codes_set = set()

    # Add revetement codes if provided explicitly
    if revetement_codes:
        coating_codes_set.update(revetement_codes)

    # Build custom products lookup
    custom_lookup = {}
    if custom_products:
        for cp in custom_products:
            custom_lookup[cp['code'].upper()] = cp

    # First pass: group products by family
    product_groups = OrderedDict()
    for family in FAMILY_ORDER:
        product_groups[family] = []

    # Process each product code
    for code in product_codes:
        # Check custom products first (SP articles)
        if code.upper() in custom_lookup:
            product = custom_lookup[code.upper()]
            family = product.get('famille', '').lower()
            if family not in product_groups:
                product_groups[family] = []
            product_groups[family].append(product)
            continue

        # Look up all pages for this product (multi-page support)
        pages = find_product_pages(code, references_dir)

        # Try with base code if coating suffix present
        if not pages and '-' in code:
            base_code = '-'.join(code.split('-')[:-1])
            pages = find_product_pages(base_code, references_dir)
            if pages:
                potential_coating = code.split('-')[-1]
                if len(potential_coating) <= 3:
                    coating_codes_set.add(potential_coating)

        if not pages:
            skipped.append({'code': code, 'reason': 'Product not found in references'})
            continue

        # Get family from first page
        family = pages[0].get('famille', '').lower()

        # Skip fiches-existantes (Phase 5 handles these)
        if family == 'fiches-existantes':
            skipped.append({'code': code, 'reason': 'Fiches-existantes handled in Phase 5'})
            continue

        # Check if family has a layout mapping
        if family not in FAMILY_LAYOUT_MAP:
            skipped.append({'code': code, 'reason': f'No layout mapping for family: {family}'})
            continue

        # Add to product groups (all pages for multi-page products)
        for product in pages:
            # Check if product has coating information in dimensions
            for dimension in product.get('dimensions', []):
                dim_name = dimension.get('name', '').lower()
                if 'revetement' in dim_name or 'revêtement' in dim_name:
                    valeur = dimension.get('valeur', '')
                    for potential_code in valeur.split(','):
                        cleaned = potential_code.strip()
                        if cleaned and len(cleaned) <= 3 and cleaned.isupper():
                            coating_codes_set.add(cleaned)

            product_groups[family].append(product)

    # Assemble complete document with cover, TOC, separators, and product slides
    if devis_info is None:
        devis_info = {}

    assembly_result = assemble_document(
        product_groups=product_groups,
        prs=prs,
        devis_info=devis_info,
        references_dir=references_dir,
        project_root=project_root,
        logo_path=project_root / 'Delagrave' / 'images' / 'logo_delagrave_emsm.png'
    )
    all_warnings = assembly_result.get('warnings', [])

    # Generate coating slides (add at end of document)
    logo_path = project_root / 'Delagrave' / 'images' / 'logo_delagrave_emsm.png'
    revetements_added = []
    if coating_codes_set:
        try:
            revetements_added = _add_revetement_slides(
                prs,
                list(coating_codes_set),
                references_dir,
                project_root,
                logo_path
            )
        except Exception as e:
            all_warnings.append({"code": "REVETEMENTS", "message": f"Erreur revetements: {str(e)}"})

    # Add page numbers LAST (after all slides including revetements)
    from gendoc.generators.modern_template import add_page_numbers
    add_page_numbers(prs)

    # Ensure output directory exists
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Save presentation
    prs.save(str(output_path))

    return {
        'slides_generated': assembly_result['product_count'],
        'total_pages': assembly_result['total_pages'] + len(revetements_added),
        'revetements_added': revetements_added,
        'skipped': skipped,
        'warnings': all_warnings
    }
