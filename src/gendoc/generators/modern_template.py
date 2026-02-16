"""
Modern template for Delagrave tech sheets — "Technical Showcase" design.

Palette from delagrave.fr (#58BAE1 accent, #1E73BE brand).
Layout: colored title band, image-left/text-right, asymmetric separators.
All slides built from blank layout with code-positioned elements.
"""

from pathlib import Path
from typing import Dict, List, Any, Optional

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE


# ──────────────────────────────────────────────────────────────
# Delagrave Color Palette (from delagrave.fr)
# ──────────────────────────────────────────────────────────────
ACCENT = RGBColor(0x58, 0xBA, 0xE1)         # #58BAE1 - primary accent
BRAND = RGBColor(0x1E, 0x73, 0xBE)          # #1E73BE - titles, headers
TEXT_MAIN = RGBColor(0x33, 0x33, 0x33)       # #333333 - body text
TEXT_LIGHT = RGBColor(0x76, 0x76, 0x76)      # #767676 - secondary text
BG_CARD = RGBColor(0xF2, 0xF4, 0xF7)        # #F2F4F7 - card backgrounds
BG_ALT = RGBColor(0xF8, 0xF9, 0xFB)         # #F8F9FB - alternating rows
DIVIDER = RGBColor(0xD8, 0xD8, 0xD8)        # #D8D8D8 - borders
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = 'Calibri'
BULLET = "\u25B8"  # ▸ small right-pointing triangle

# Default company info (overridden by _entreprise.md if found)
_COMPANY_INFO = {
    'nom': 'SAS DELAGRAVE EMSM',
    'adresse': '350, Rue Blingue - 27610 Romilly sur Andelle',
    'site_web': 'www.delagrave.fr',
    'legal_1': 'S.A.S au capital de 300 000 \u20ac \u2013 RCS EVREUX B 879 414 803 \u2013 SIRET 879 414 803 00012 \u2013 NAF 7112B \u2013 NII : FR 09 879 414 803',
    'legal_2': 'N\u00b0 enregistrement art. L. 541-10 du code de l\u2019Environnement : FR029760_10KNMW',
    'tagline': 'DELAGRAVE EMSM \u2014 Fabricant fran\u00e7ais de mobilier de laboratoire',
}


def _load_company_info(references_dir: Optional[Path] = None) -> dict:
    """Load company info from _entreprise.md, falling back to defaults."""
    info = dict(_COMPANY_INFO)
    if references_dir is None:
        return info
    filepath = references_dir / '_entreprise.md'
    if not filepath.exists():
        return info
    try:
        text = filepath.read_text(encoding='utf-8')
        for line in text.split('\n'):
            line = line.strip()
            if line.startswith('|') and not line.startswith('| Cle') and not line.startswith('|--'):
                parts = [p.strip() for p in line.split('|')]
                parts = [p for p in parts if p]
                if len(parts) >= 2:
                    info[parts[0]] = parts[1]
    except Exception:
        pass
    return info


# Module-level cache (loaded once per generation via init_company_info)
_company = dict(_COMPANY_INFO)


def init_company_info(references_dir: Path):
    """Initialize company info from _entreprise.md. Call once before generating."""
    global _company
    _company = _load_company_info(references_dir)


# ──────────────────────────────────────────────────────────────
# Primitive Helpers
# ──────────────────────────────────────────────────────────────

def _rect(slide, x, y, w, h, color):
    """Solid rectangle, no border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_shadow(shape):
    """Add subtle outer drop shadow to a shape via XML."""
    from pptx.oxml import parse_xml
    effect_xml = (
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:outerShdw blurRad="76200" dist="38100" dir="5400000" algn="t" rotWithShape="0">'
        '<a:srgbClr val="000000"><a:alpha val="23000"/></a:srgbClr>'
        '</a:outerShdw>'
        '</a:effectLst>'
    )
    effectLst = parse_xml(effect_xml)
    shape._element.spPr.append(effectLst)


def _hline(slide, x, y, w, color=None):
    """Horizontal thin line."""
    _rect(slide, x, y, w, Cm(0.04), color or DIVIDER)


# ──────────────────────────────────────────────────────────────
# Header / Footer / Title Band
# ──────────────────────────────────────────────────────────────

def _header(slide, prs, logo_path):
    """Slim header: logo left + thin accent line."""
    sw = prs.slide_width
    if logo_path and logo_path.exists():
        logo_dir = logo_path.parent
        logo_file = logo_dir / 'delagrave2022_high.png'
        if not logo_file.exists():
            logo_file = logo_path
        slide.shapes.add_picture(str(logo_file), Cm(0.6), Cm(0.25), height=Cm(1.1))
    # Full-width accent line
    _rect(slide, 0, Cm(1.6), sw, Cm(0.06), ACCENT)


def _footer(slide, prs):
    """Footer: accent line + company info + legal lines at bottom of slide."""
    sw = prs.slide_width
    # Accent line
    _rect(slide, 0, Cm(28.0), sw, Cm(0.05), ACCENT)
    # Website (left)
    box = slide.shapes.add_textbox(Cm(0.6), Cm(28.15), Cm(4), Cm(0.4))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = _company['site_web']
    r.font.size = Pt(7)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = BRAND
    # Company name + address (centered)
    box = slide.shapes.add_textbox(Cm(0.4), Cm(28.15), sw - Cm(0.8), Cm(0.4))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"{_company['nom']} - {_company['adresse']}"
    r.font.size = Pt(7)
    r.font.name = FONT
    r.font.color.rgb = TEXT_LIGHT
    # Legal info (two lines, centered)
    box = slide.shapes.add_textbox(Cm(0.4), Cm(28.55), sw - Cm(0.8), Cm(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r = p1.add_run()
    r.text = _company['legal_1']
    r.font.size = Pt(5.5)
    r.font.name = FONT
    r.font.color.rgb = TEXT_LIGHT
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run()
    r.text = _company['legal_2']
    r.font.size = Pt(5.5)
    r.font.name = FONT
    r.font.color.rgb = TEXT_LIGHT


def _title_band(slide, prs, text):
    """Full-width colored title band with white text."""
    sw = prs.slide_width
    _rect(slide, 0, Cm(1.85), sw, Cm(1.3), BRAND)
    box = slide.shapes.add_textbox(Cm(0.8), Cm(1.88), sw - Cm(1.6), Cm(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = WHITE


def _section_label(slide, text, x, y, w):
    """Section header: uppercase brand label + short accent underline."""
    box = slide.shapes.add_textbox(x, y, w, Cm(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = BRAND
    _rect(slide, x, y + Cm(0.48), min(Cm(2.5), w), Cm(0.05), ACCENT)
    return y + Cm(0.65)


# ──────────────────────────────────────────────────────────────
# Content Builders
# ──────────────────────────────────────────────────────────────

def _text_with_bullets(slide, text, x, y, w, h):
    """Text with ▸ blue triangle bullets, justified, hanging indent."""
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    if not lines:
        return y

    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.text = ""

    indent = Cm(0.4)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        bullet_run = p.add_run()
        bullet_run.text = f"{BULLET} "
        bullet_run.font.color.rgb = ACCENT
        bullet_run.font.size = Pt(10)
        bullet_run.font.name = FONT
        bullet_run.font.bold = True

        text_run = p.add_run()
        text_run.text = line
        text_run.font.size = Pt(9)
        text_run.font.name = FONT
        text_run.font.color.rgb = TEXT_MAIN

        p.alignment = PP_ALIGN.JUSTIFY
        p.space_after = Pt(3)
        pPr = p._p.get_or_add_pPr()
        pPr.set('marL', str(indent))
        pPr.set('indent', str(-indent))

    return y + h


def _dimensions_table(slide, dims, x, y, w):
    """Dimensions table: brand header + alternating data rows."""
    if not dims:
        return y

    row_h = Cm(0.55)

    # Header row
    _rect(slide, x, y, w, row_h, BRAND)
    hb = slide.shapes.add_textbox(x + Cm(0.3), y, int(w * 0.6), row_h)
    tf = hb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Caract\u00e9ristique"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = WHITE

    hb2 = slide.shapes.add_textbox(x + int(w * 0.6), y, int(w * 0.37), row_h)
    tf = hb2.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.text = "Valeur"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = WHITE

    y += row_h
    count = 0

    for dim in dims:
        name = dim.get('name', '').strip()
        valeur = dim.get('valeur', '').strip()
        prefix = dim.get('prefix', '').strip()
        if prefix:
            valeur = f"{prefix}{valeur}"
        if not name or not valeur:
            continue

        ry = y + count * row_h
        bg = BG_ALT if count % 2 == 0 else WHITE
        _rect(slide, x, ry, w, row_h, bg)

        lb = slide.shapes.add_textbox(x + Cm(0.3), ry, int(w * 0.58), row_h)
        tf = lb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(8)
        p.font.name = FONT
        p.font.color.rgb = TEXT_MAIN

        vb = slide.shapes.add_textbox(x + int(w * 0.58), ry, int(w * 0.37), row_h)
        tf = vb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.text = valeur
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = BRAND

        count += 1

    _rect(slide, x, y + count * row_h, w, Cm(0.04), ACCENT)
    return y + count * row_h + Cm(0.15)


# ──────────────────────────────────────────────────────────────
# Image Insertion
# ──────────────────────────────────────────────────────────────

def _insert_image(slide, images, project_root, x, y, max_w, max_h):
    """Insert first valid image, fitted and centered in area."""
    for img_data in images:
        chemin = img_data.get('chemin', '').strip()
        if not chemin or chemin.endswith('.missing'):
            continue
        path = project_root / chemin
        if not path.exists():
            continue
        try:
            from PIL import Image as PILImage
            with PILImage.open(str(path)) as im:
                iw, ih = im.size
        except Exception:
            continue
        if iw <= 0:
            continue

        ratio = ih / iw
        fit_w = max_w
        fit_h = int(fit_w * ratio)
        if fit_h > max_h:
            fit_h = max_h
            fit_w = int(fit_h / ratio)

        cx = x + (max_w - fit_w) // 2
        cy = y + (max_h - fit_h) // 2

        try:
            slide.shapes.add_picture(str(path), cx, cy, fit_w, fit_h)
            return True
        except Exception:
            continue
    return False


def _insert_all_images(slide, images, project_root, x, y, max_w, max_h_each):
    """Insert all valid images, stacked vertically."""
    cur_y = y
    for img_data in images:
        chemin = img_data.get('chemin', '').strip()
        if not chemin or chemin.endswith('.missing'):
            continue
        path = project_root / chemin
        if not path.exists():
            continue
        try:
            from PIL import Image as PILImage
            with PILImage.open(str(path)) as im:
                iw, ih = im.size
        except Exception:
            continue
        if iw <= 0:
            continue

        ratio = ih / iw
        fit_w = max_w
        fit_h = int(fit_w * ratio)
        if fit_h > max_h_each:
            fit_h = max_h_each
            fit_w = int(fit_h / ratio)

        cx = x + (max_w - fit_w) // 2
        try:
            slide.shapes.add_picture(str(path), cx, cur_y, fit_w, fit_h)
            cur_y += fit_h + Cm(0.3)
        except Exception:
            continue


# ──────────────────────────────────────────────────────────────
# Page Builders
# ──────────────────────────────────────────────────────────────

def build_cover(prs, devis_info, logo_path):
    """Cover: logo, accent line, title, info card, branded bottom band."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    sw = prs.slide_width
    sh = prs.slide_height

    # Logo centered at top
    if logo_path and logo_path.exists():
        logo_dir = logo_path.parent
        logo_file = logo_dir / 'delagrave2022_high.png'
        if not logo_file.exists():
            logo_file = logo_path
        # Center logo horizontally
        from PIL import Image as PILImage
        try:
            with PILImage.open(str(logo_file)) as im:
                ratio = im.size[0] / im.size[1]
            logo_h = Cm(2.5)
            logo_w = int(logo_h * ratio)
            slide.shapes.add_picture(
                str(logo_file), (sw - logo_w) // 2, Cm(2), width=logo_w, height=logo_h
            )
        except Exception:
            slide.shapes.add_picture(str(logo_file), Cm(1.5), Cm(2), height=Cm(2.5))

    # Accent line
    _rect(slide, Cm(3), Cm(5.5), sw - Cm(6), Cm(0.06), ACCENT)

    # "FICHES TECHNIQUES"
    box = slide.shapes.add_textbox(Cm(1), Cm(7), sw - Cm(2), Cm(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "FICHES TECHNIQUES"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = ACCENT
    rPr = r._r.get_or_add_rPr()
    rPr.set('spc', '400')

    # Project title
    title_text = devis_info.get(
        'titre_affaire',
        devis_info.get('client', 'Dossier Technique')
    )
    box = slide.shapes.add_textbox(Cm(2), Cm(10), sw - Cm(4), Cm(3))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = BRAND

    # Info card (centered, with subtle background)
    info_lines = []
    if devis_info.get('numero_devis'):
        info_lines.append(f"Devis N\u00b0 {devis_info['numero_devis']}")
    if devis_info.get('date'):
        info_lines.append(f"Date : {devis_info['date']}")
    if devis_info.get('client'):
        info_lines.append(f"Client : {devis_info['client']}")

    if info_lines:
        card_w = Cm(12)
        card_x = (sw - card_w) // 2
        card_y = Cm(14.5)
        card_h = Cm(3.5)
        _rect(slide, card_x, card_y, card_w, card_h, BG_CARD)
        # Left accent bar on card
        _rect(slide, card_x, card_y, Cm(0.12), card_h, ACCENT)

        box = slide.shapes.add_textbox(card_x + Cm(0.8), card_y + Cm(0.4), card_w - Cm(1.2), card_h - Cm(0.8))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(info_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.text = line
            p.font.size = Pt(13)
            p.font.name = FONT
            p.font.color.rgb = TEXT_MAIN
            p.space_after = Pt(6)

    # Bottom band
    band_h = Cm(2)
    _rect(slide, 0, sh - band_h, sw, band_h, BRAND)
    _rect(slide, 0, sh - band_h, sw, Cm(0.06), ACCENT)
    box = slide.shapes.add_textbox(Cm(1), sh - band_h + Cm(0.55), sw - Cm(2), Cm(0.8))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = _company['tagline']
    r.font.size = Pt(11)
    r.font.name = FONT
    r.font.color.rgb = WHITE
    r.font.bold = True


# ── TOC layout constants ──
_TOC_MAX_Y = Cm(27.5)              # Bottom limit before footer zone
_TOC_FIRST_START_Y = Cm(4.2)       # After branded header on first page
_TOC_CONT_START_Y = Cm(2.8)        # After smaller header on continuation pages
_TOC_FAMILY_H = Cm(0.7)            # Height of family header row
_TOC_ROW_H = Cm(0.5)               # Height of product row
_TOC_FAMILY_GAP = Cm(0.35)         # Gap after last product in family


def estimate_toc_pages(toc_entries):
    """Pre-calculate how many TOC slides are needed for page numbering."""
    y = _TOC_FIRST_START_Y
    pages = 1

    for entry in toc_entries:
        if not entry.get('products'):
            continue

        # Need room for family header + at least 1 product row
        if y + _TOC_FAMILY_H + _TOC_ROW_H > _TOC_MAX_Y:
            pages += 1
            y = _TOC_CONT_START_Y

        y += _TOC_FAMILY_H

        for _product in entry['products']:
            if y + _TOC_ROW_H > _TOC_MAX_Y:
                pages += 1
                y = _TOC_CONT_START_Y
            y += _TOC_ROW_H

        y += _TOC_FAMILY_GAP

    return pages


def _toc_new_slide(prs, sw, is_first):
    """Create a TOC slide with appropriate header. Returns (slide, start_y)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    if is_first:
        # Full branded header
        _rect(slide, 0, 0, sw, Cm(3.2), BRAND)
        _rect(slide, 0, Cm(3.2), sw, Cm(0.08), ACCENT)

        box = slide.shapes.add_textbox(Cm(1.5), Cm(0.8), sw - Cm(3), Cm(1.8))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "SOMMAIRE"
        r.font.size = Pt(26)
        r.font.bold = True
        r.font.name = FONT
        r.font.color.rgb = WHITE
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', '250')

        return slide, _TOC_FIRST_START_Y
    else:
        # Smaller continuation header
        _rect(slide, 0, 0, sw, Cm(2.0), BRAND)
        _rect(slide, 0, Cm(2.0), sw, Cm(0.08), ACCENT)

        box = slide.shapes.add_textbox(Cm(1.5), Cm(0.4), sw - Cm(3), Cm(1.2))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "SOMMAIRE (suite)"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.name = FONT
        r.font.color.rgb = WHITE
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', '250')

        return slide, _TOC_CONT_START_Y


def build_toc(prs, toc_entries):
    """TOC with automatic page breaks when content overflows."""
    sw = prs.slide_width

    mx = Cm(1.5)
    cw = sw - Cm(3)

    slide, y = _toc_new_slide(prs, sw, is_first=True)

    for entry in toc_entries:
        if not entry['products']:
            continue

        # Ensure room for family header + at least 1 product
        if y + _TOC_FAMILY_H + _TOC_ROW_H > _TOC_MAX_Y:
            slide, y = _toc_new_slide(prs, sw, is_first=False)

        # Family header with accent dot
        _rect(slide, mx, y + Cm(0.12), Cm(0.22), Cm(0.22), ACCENT)
        box = slide.shapes.add_textbox(mx + Cm(0.45), y, cw - Cm(0.45), Cm(0.55))
        p = box.text_frame.paragraphs[0]
        p.text = entry['family_display']
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = BRAND
        y += _TOC_FAMILY_H

        # Product rows
        for i, product in enumerate(entry['products']):
            if y + _TOC_ROW_H > _TOC_MAX_Y:
                slide, y = _toc_new_slide(prs, sw, is_first=False)

            bg = BG_ALT if i % 2 == 0 else WHITE
            _rect(slide, mx + Cm(0.45), y, cw - Cm(0.45), _TOC_ROW_H, bg)

            titre = product['titre']
            if len(titre) > 55:
                titre = titre[:52] + '...'

            box = slide.shapes.add_textbox(mx + Cm(0.8), y, cw - Cm(3), _TOC_ROW_H)
            tf = box.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = f"{product['code']} \u2014 {titre}"
            p.font.size = Pt(9)
            p.font.name = FONT
            p.font.color.rgb = TEXT_MAIN

            box = slide.shapes.add_textbox(mx + cw - Cm(1.8), y, Cm(1.5), _TOC_ROW_H)
            tf = box.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            p.text = str(product['page_number'])
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.name = FONT
            p.font.color.rgb = BRAND

            y += _TOC_ROW_H

        y += _TOC_FAMILY_GAP


def build_separator(prs, family_name, family_display):
    """Asymmetric separator: left brand band + white area."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    sw = prs.slide_width
    sh = prs.slide_height

    # Left vertical band (35% width, full height)
    band_w = int(sw * 0.35)
    _rect(slide, 0, 0, band_w, sh, BRAND)

    # Accent line at right edge of band
    _rect(slide, band_w, 0, Cm(0.08), sh, ACCENT)

    # Family name in the white area, vertically centered
    text_x = band_w + Cm(1.5)
    text_w = sw - band_w - Cm(2.5)
    box = slide.shapes.add_textbox(text_x, Cm(10), text_w, Cm(6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.text = family_display
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = BRAND

    # Short accent line below name
    _rect(slide, text_x, Cm(16.5), Cm(4), Cm(0.08), ACCENT)

    # Small section number or decorative element in the band
    box = slide.shapes.add_textbox(Cm(1), Cm(12), band_w - Cm(2), Cm(3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "\u25B8"
    r.font.size = Pt(48)
    r.font.color.rgb = ACCENT
    r.font.name = FONT


# ──────────────────────────────────────────────────────────────
# Product Slide Builders
# ──────────────────────────────────────────────────────────────

def build_product_slide(prs, product, family, project_root, logo_path):
    """Dispatch to family-specific builder. Returns list of warning strings."""
    try:
        if family in ('armoire-securite', 'enceinte-ventilee'):
            # Both armoire-securite and enceinte-ventilee use Option C 2-page template
            # armoire-securite: commit 0b3600b (2026-02-15)
            # enceinte-ventilee: commit 0cee8d5 (2026-02-16)
            return _build_armoire_slide(prs, product, project_root, logo_path)
        elif family in ('equipement', 'elec-sorb', 'complements'):
            return _build_simple_slide(prs, product, project_root, logo_path)
        elif family == 'revetement':
            return _build_revetement_slide(prs, product, project_root, logo_path)
        else:
            return _build_standard_slide(prs, product, family, project_root, logo_path)
    except Exception as e:
        return [f"Erreur inattendue sur slide: {str(e)}"]


def _build_standard_slide(prs, product, family, project_root, logo_path):
    """
    Standard slide: image LEFT on card + text RIGHT.

    Layout zones (no overlap guaranteed):
    - Header:    0     - 1.65cm  (logo + accent line)
    - Title:     1.85  - 3.15cm  (brand band with white text)
    - Image:     3.5   - 27cm    (x: 0.5 - 9cm, card background)
    - Text:      3.5   - 27cm    (x: 9.8 - 20.3cm)
    - Footer:    27.3  - 28cm    (divider + info)

    Returns list of warning strings.
    """
    try:
        warnings = []
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        _header(slide, prs, logo_path)
        _title_band(slide, prs, product.get('titre', ''))
        _footer(slide, prs)

        # Image zone (left) — card background
        img_x = Cm(0.5)
        img_w = Cm(8.5)
        img_top = Cm(3.5)
        img_bot = Cm(27)
        img_h = img_bot - img_top

        # White card with drop shadow for image
        card = _rect(slide, img_x, img_top, img_w, img_h, WHITE)
        card.line.color.rgb = DIVIDER
        card.line.width = Pt(0.5)
        _add_shadow(card)
        # Accent line on left edge of card
        _rect(slide, img_x, img_top, Cm(0.08), img_h, ACCENT)

        # Image centered inside white card
        images = product.get('images', [])
        inserted = _insert_image(
            slide, images, project_root,
            img_x + Cm(0.5), img_top + Cm(0.8), img_w - Cm(1.0), img_h - Cm(1.6)
        )
        if not inserted and images:
            # Had images but none could be inserted
            code = product.get('code', '?')
            warnings.append(f"Aucune image inseree pour {code}")

        # Text zone (right)
        tx = Cm(9.8)
        tw = Cm(10.5)
        y = Cm(3.5)

        # Separate revetement info from dimensions
        actual_dims = []
        revetement_text = ''
        no_dims_families = {'tables-en'}
        for dim in product.get('dimensions', []):
            name_lower = dim.get('name', '').lower()
            if 'revetement' in name_lower or 'rev\u00eatement' in name_lower:
                revetement_text = dim.get("valeur", "")
            elif 'page' in name_lower:
                continue
            else:
                actual_dims.append(dim)

        # DESCRIPTION
        texte = product.get('texte', '').strip()
        if texte and texte.lower() != 'aucune':
            y = _section_label(slide, 'Description', tx, y, tw)
            raw_lines = [l for l in texte.split('\n') if l.strip()]
            visual_lines = sum(max(1, len(l) // 45 + 1) for l in raw_lines)
            desc_h = min(Cm(15), max(Cm(3), Cm(0.42) * visual_lines + Cm(1)))
            y = _text_with_bullets(slide, texte, tx, y, tw, desc_h)
            y += Cm(0.5)

        # DIMENSIONS (skip for tables-en)
        if actual_dims and family not in no_dims_families:
            y = _section_label(slide, 'Dimensions', tx, y, tw)
            y = _dimensions_table(slide, actual_dims, tx, y, tw)
            y += Cm(0.5)

        # REVETEMENTS
        if revetement_text:
            y = _section_label(slide, 'Rev\u00eatements disponibles', tx, y, tw)
            box = slide.shapes.add_textbox(tx + Cm(0.2), y, tw - Cm(0.2), Cm(1.5))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = revetement_text
            p.font.size = Pt(9)
            p.font.name = FONT
            p.font.color.rgb = TEXT_MAIN

        return warnings
    except Exception as e:
        return [f"Erreur inattendue sur slide: {str(e)}"]


def _build_revetement_slide(prs, product, project_root, logo_path):
    """Revetement slide: full-width sections + bottom split for images.

    Returns list of warning strings.
    """
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        _header(slide, prs, logo_path)
        _title_band(slide, prs, product.get('titre', ''))
        _footer(slide, prs)

        fw = prs.slide_width - Cm(1.2)
        lx = Cm(0.6)
        y = Cm(3.5)

        texte = product.get('texte', '')
        blocks = [b.strip() for b in texte.split('\n\n') if b.strip()]

        description = blocks[0] if len(blocks) >= 1 else ''
        mise_en_oeuvre = blocks[1] if len(blocks) >= 2 else ''
        finition = '\n\n'.join(blocks[2:]) if len(blocks) >= 3 else ''

        applications = ''
        for dim in product.get('dimensions', []):
            name_lower = dim.get('name', '').lower()
            if 'application' in name_lower:
                applications = dim.get('valeur', '')

        # APPLICATIONS (full width)
        if applications:
            y = _section_label(slide, 'Applications', lx, y, fw)
            box = slide.shapes.add_textbox(lx + Cm(0.2), y, fw - Cm(0.2), Cm(1.2))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = applications
            p.font.size = Pt(10)
            p.font.name = FONT
            p.font.color.rgb = TEXT_MAIN
            y += Cm(1.5)

        # CARACTERISTIQUES (full width)
        if description:
            y = _section_label(slide, 'Caract\u00e9ristiques', lx, y, fw)
            lines_count = len([l for l in description.split('\n') if l.strip()])
            desc_h = min(Cm(8), max(Cm(2), Cm(0.45) * lines_count + Cm(0.5)))
            y = _text_with_bullets(slide, description, lx, y, fw, desc_h)
            y += Cm(0.4)

        # Split: text left + images right
        text_w = Cm(11)
        img_x = Cm(12.5)
        img_w = Cm(7.5)

        if mise_en_oeuvre:
            y = _section_label(slide, 'Mise en \u0153uvre', lx, y, text_w)
            lines_count = len([l for l in mise_en_oeuvre.split('\n') if l.strip()])
            h = min(Cm(5), max(Cm(1.5), Cm(0.45) * lines_count + Cm(0.5)))
            y = _text_with_bullets(slide, mise_en_oeuvre, lx, y, text_w, h)
            y += Cm(0.3)

        if finition:
            y = _section_label(slide, 'Finition', lx, y, text_w)
            lines_count = len([l for l in finition.split('\n') if l.strip()])
            h = min(Cm(4), max(Cm(1), Cm(0.45) * lines_count + Cm(0.5)))
            _text_with_bullets(slide, finition, lx, y, text_w, h)

        # Images (right column)
        _insert_all_images(
            slide, product.get('images', []), project_root,
            img_x, Cm(12), img_w, Cm(7)
        )
        return []
    except Exception as e:
        return [f"Erreur inattendue sur slide: {str(e)}"]


def _build_simple_slide(prs, product, project_root, logo_path):
    """Simple slide for equipement/elec-sorb/complements: one slide per image.

    Multi-image products (e.g. BC1Vx with 5 datasheets) get one slide per image,
    matching original VBA behavior. Single-image products are unchanged.

    Returns list of warning strings.
    """
    try:
        warnings = []
        images = [img for img in product.get('images', [])
                  if img.get('chemin', '').strip()
                  and not img.get('chemin', '').strip().endswith('.missing')]

        if not images:
            images = [{}]  # Still create one slide even with no images

        sw = prs.slide_width
        for img_data in images:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            _header(slide, prs, logo_path)
            _title_band(slide, prs, product.get('titre', ''))
            _footer(slide, prs)

            chemin = img_data.get('chemin', '').strip()
            if chemin:
                inserted = _insert_image(
                    slide, [img_data], project_root,
                    Cm(1.5), Cm(3.8), sw - Cm(3), Cm(22)
                )
                if not inserted:
                    code = product.get('code', '?')
                    warnings.append(f"Image non inseree pour {code}: {chemin}")

        return warnings
    except Exception as e:
        return [f"Erreur inattendue sur slide: {str(e)}"]


# ──────────────────────────────────────────────────────────────
# Page Numbering
# ──────────────────────────────────────────────────────────────

def add_page_numbers(prs):
    """Page numbers top-right of footer zone."""
    sw = prs.slide_width
    for i, slide in enumerate(prs.slides):
        box = slide.shapes.add_textbox(sw - Cm(1.8), Cm(28.12), Cm(1.2), Cm(0.4))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = str(i + 1)
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.name = FONT
        r.font.color.rgb = BRAND


# ═════════════════════════════════════════════════════════════
# Armoire-securite: Option C Template (2-page layout)
# ═════════════════════════════════════════════════════════════
# Security cabinets require dedicated layout for:
#   - Dual images (product photo + interior schema)
#   - Certification badges (EN 14470-1, FM, etc.)
#   - Construction details (keyword-based bullets)
#   - Full-width specifications table
#
# This is the only family using multi-page template.
# Added: commit 0b3600b (2026-02-15)
# ═════════════════════════════════════════════════════════════

def _parse_armoire_texte(texte):
    """Parse armoire texte field into description, certificats, and fonction.

    Format expected in MD texte field:
        Description text (multiple lines)
        ---CERTIFICATS---
        Cert line 1
        Cert line 2
        ---FONCTION---
        Keyword : description
        Keyword : description

    Args:
        texte (str): Raw text from product['texte'] field

    Returns:
        tuple: (description: str, certificats: List[str], fonctions: List[tuple|str])
               fonctions are tuples of (keyword, description) or strings if no colon

    Note:
        Added in commit 0b3600b for armoire-securite family (Option C template).
        This parsing enables structured presentation of security cabinet specifications.
    """
    description = texte
    certificats = []
    fonctions = []

    if '---CERTIFICATS---' in texte:
        parts = texte.split('---CERTIFICATS---')
        description = parts[0].strip()
        rest = parts[1]

        if '---FONCTION---' in rest:
            cert_part, fonc_part = rest.split('---FONCTION---')
        else:
            cert_part = rest
            fonc_part = ''

        certificats = [l.strip() for l in cert_part.strip().split('\n') if l.strip()]
        for line in fonc_part.strip().split('\n'):
            line = line.strip()
            if not line:
                continue
            if ' : ' in line:
                kw, desc = line.split(' : ', 1)
                fonctions.append((kw.strip(), desc.strip()))
            else:
                fonctions.append(line)

    return description, certificats, fonctions


def _ref_line(slide, prs, ref_text):
    """Small reference line right-aligned just below the title band."""
    if not ref_text:
        return
    sw = prs.slide_width
    box = slide.shapes.add_textbox(sw - Cm(10.5), Cm(3.2), Cm(10), Cm(0.4))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = ref_text
    r.font.size = Pt(8)
    r.font.name = FONT
    r.font.color.rgb = TEXT_LIGHT
    r.font.italic = True


def _rich_bullets(slide, items, x, y, w, h):
    """Bullet list with bold keyword + normal description.
    items: list of (keyword, description) tuples or plain strings.
    """
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.text = ""
    indent = Cm(0.4)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        bullet_run = p.add_run()
        bullet_run.text = f"{BULLET} "
        bullet_run.font.color.rgb = ACCENT
        bullet_run.font.size = Pt(10)
        bullet_run.font.name = FONT
        bullet_run.font.bold = True

        if isinstance(item, tuple):
            keyword, desc = item
            kw_run = p.add_run()
            kw_run.text = f"{keyword} : "
            kw_run.font.size = Pt(9)
            kw_run.font.name = FONT
            kw_run.font.color.rgb = BRAND
            kw_run.font.bold = True
            desc_run = p.add_run()
            desc_run.text = desc
            desc_run.font.size = Pt(9)
            desc_run.font.name = FONT
            desc_run.font.color.rgb = TEXT_MAIN
        else:
            text_run = p.add_run()
            text_run.text = item
            text_run.font.size = Pt(9)
            text_run.font.name = FONT
            text_run.font.color.rgb = TEXT_MAIN

        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        pPr = p._p.get_or_add_pPr()
        pPr.set('marL', str(indent))
        pPr.set('indent', str(-indent))

    return y + h


def _certif_badges(slide, certificats, x, y, w):
    """Compact certification list."""
    box = slide.shapes.add_textbox(x, y, w, Cm(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.text = ""

    for i, cert in enumerate(certificats):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet_run = p.add_run()
        bullet_run.text = f"{BULLET} "
        bullet_run.font.color.rgb = ACCENT
        bullet_run.font.size = Pt(8)
        bullet_run.font.name = FONT
        bullet_run.font.bold = True

        text_run = p.add_run()
        text_run.text = cert
        text_run.font.size = Pt(8)
        text_run.font.name = FONT
        text_run.font.color.rgb = TEXT_MAIN

        p.space_after = Pt(1)
        p.alignment = PP_ALIGN.LEFT

    return y + Cm(2.5)


def _wide_dimensions_table(slide, dims, x, y, w):
    """Wide dimensions table. Two side-by-side columns if > 6 rows."""
    if not dims:
        return y

    row_h = Cm(0.50)

    if len(dims) > 6:
        mid = (len(dims) + 1) // 2
        col1 = dims[:mid]
        col2 = dims[mid:]
        col_w = (w - Cm(0.5)) / 2

        _draw_dim_column(slide, col1, x, y, col_w, row_h)
        _draw_dim_column(slide, col2, x + col_w + Cm(0.5), y, col_w, row_h)

        max_rows = max(len(col1), len(col2))
        return y + (max_rows + 1) * row_h + Cm(0.15)
    else:
        _draw_dim_column(slide, dims, x, y, w, row_h)
        return y + (len(dims) + 1) * row_h + Cm(0.15)


def _draw_dim_column(slide, dims, x, y, w, row_h):
    """Draw a single dimensions column with header."""
    _rect(slide, x, y, w, row_h, BRAND)
    hb = slide.shapes.add_textbox(x + Cm(0.3), y, int(w * 0.6), row_h)
    tf = hb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Caractéristique"
    p.font.size = Pt(7.5)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = WHITE

    hb2 = slide.shapes.add_textbox(x + int(w * 0.6), y, int(w * 0.37), row_h)
    tf = hb2.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.text = "Valeur"
    p.font.size = Pt(7.5)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = WHITE

    for i, dim in enumerate(dims):
        name = dim.get('name', '').strip()
        valeur = dim.get('valeur', '').strip()
        if not name or not valeur:
            continue

        ry = y + (i + 1) * row_h
        bg = BG_ALT if i % 2 == 0 else WHITE
        _rect(slide, x, ry, w, row_h, bg)

        lb = slide.shapes.add_textbox(x + Cm(0.3), ry, int(w * 0.58), row_h)
        tf = lb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(7.5)
        p.font.name = FONT
        p.font.color.rgb = TEXT_MAIN

        vb = slide.shapes.add_textbox(x + int(w * 0.58), ry, int(w * 0.37), row_h)
        tf = vb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.text = valeur
        p.font.size = Pt(7.5)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = BRAND

    _rect(slide, x, y + (len(dims) + 1) * row_h, w, Cm(0.04), ACCENT)


def _build_armoire_slide(prs, product, project_root, logo_path):
    """Build armoire-securite slides using Option C template (2 pages per product).

    armoire-securite uses a dedicated 2-page layout instead of standard single-page template:

    Page 1 - Product Presentation:
      - Left column (8.5cm card):
        * Upper half: Product photo
        * Lower half: Interior schema (if available)
      - Right column (10.5cm):
        * Description section (with bullets)
        * Certificats & Normes badges
        * Fonction / Construction (rich bullets with keywords)

    Page 2 - Technical Specifications:
      - Full-width specifications table (from dimensions data)
      - Technical plan/drawing (if available in images)

    Design rationale:
      - Security cabinets need space for certifications (EN 14470-1, FM, etc.)
      - Two images per product (photo + schema) require dedicated layout
      - Specifications table too wide for single-column layout

    Args:
        prs (Presentation): Target presentation object
        product (dict): Product data from parse_family_md
        project_root (Path): Base path for resolving image paths
        logo_path (Path): Company logo path for header

    Returns:
        list: Warning messages (e.g., missing images)

    Note:
        Added in commit 0b3600b. This is the ONLY family using multi-page layout.
        enceinte-ventilee uses standard single-page template despite similar content.
    """
    try:
        warnings = []
        texte = product.get('texte', '')
        description, certificats, fonctions = _parse_armoire_texte(texte)
        images = product.get('images', [])

        # ── PAGE 1: Presentation ──
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])

        _header(slide1, prs, logo_path)
        _title_band(slide1, prs, product.get('titre', ''))
        _ref_line(slide1, prs, product.get('ref', ''))
        _footer(slide1, prs)

        # Image zone (left) — card
        img_x = Cm(0.5)
        img_w = Cm(8.5)
        img_top = Cm(3.5)
        img_bot = Cm(27)
        img_h = img_bot - img_top

        card = _rect(slide1, img_x, img_top, img_w, img_h, WHITE)
        card.line.color.rgb = DIVIDER
        card.line.width = Pt(0.5)
        _add_shadow(card)
        _rect(slide1, img_x, img_top, Cm(0.08), img_h, ACCENT)

        # Photo produit in upper half
        photo_images = [img for img in images if 'photo' in img.get('position', '').lower()]
        if not photo_images:
            photo_images = [img for img in images if 'schema' not in img.get('position', '').lower()
                            and 'plan' not in img.get('position', '').lower()]
        if not photo_images:
            photo_images = images[:1]

        half_h = img_h // 2 - Cm(0.5)
        inserted = _insert_image(
            slide1, photo_images, project_root,
            img_x + Cm(0.5), img_top + Cm(0.4), img_w - Cm(1.0), half_h
        )
        if not inserted and photo_images:
            warnings.append(f"Aucune photo inseree pour {product.get('code', '?')}")

        # Schema interieur in lower half
        schema_images = [img for img in images if 'schema' in img.get('position', '').lower()]
        if schema_images:
            _insert_image(
                slide1, schema_images, project_root,
                img_x + Cm(0.5), img_top + half_h + Cm(1.0),
                img_w - Cm(1.0), half_h
            )

        # Text zone (right)
        tx = Cm(9.8)
        tw = Cm(10.5)
        y = Cm(3.5)

        if description:
            y = _section_label(slide1, 'Description', tx, y, tw)
            raw_lines = [l for l in description.split('\n') if l.strip()]
            visual_lines = sum(max(1, len(l) // 45 + 1) for l in raw_lines)
            desc_h = min(Cm(12), max(Cm(2.5), Cm(0.42) * visual_lines + Cm(0.5)))
            y = _text_with_bullets(slide1, description, tx, y, tw, desc_h)
            y += Cm(0.4)

        if certificats:
            y = _section_label(slide1, 'Certificats & Normes', tx, y, tw)
            y = _certif_badges(slide1, certificats, tx, y, tw)
            y += Cm(0.4)

        if fonctions:
            y = _section_label(slide1, 'Fonction / Construction', tx, y, tw)
            remaining_h = Cm(27) - y
            _rich_bullets(slide1, fonctions, tx, y, tw, remaining_h)

        # ── PAGE 2: Specs + plan technique ──
        slide2 = prs.slides.add_slide(prs.slide_layouts[0])
        sw = prs.slide_width

        _header(slide2, prs, logo_path)
        code = product.get('code', '')
        _title_band(slide2, prs, f"Caractéristiques techniques — {code}")
        _ref_line(slide2, prs, product.get('ref', ''))
        _footer(slide2, prs)

        y = Cm(3.8)
        dims = product.get('dimensions', [])
        table_w = sw - Cm(2)
        table_x = Cm(1.0)

        _section_label(slide2, 'Caractéristiques techniques', table_x, y, table_w)
        y += Cm(0.7)
        y = _wide_dimensions_table(slide2, dims, table_x, y, table_w)
        y += Cm(0.8)

        plan_images = [img for img in images if 'plan' in img.get('position', '').lower()]

        if plan_images:
            _section_label(slide2, 'Plan dimensionnel', table_x, y, table_w)
            y += Cm(0.7)
            remaining_h = Cm(27) - y
            _insert_image(
                slide2, plan_images, project_root,
                table_x, y, table_w, remaining_h
            )

        return warnings
    except Exception as e:
        return [f"Erreur inattendue sur slide armoire: {str(e)}"]
