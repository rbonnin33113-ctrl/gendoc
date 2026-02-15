"""
Document assembly layer for Delagrave PowerPoint generation.

This module transforms individual product slides into a complete professional
PowerPoint document with cover page, chapter separator pages, and table of contents.

This is a pure library module (no I/O beyond file operations).
"""

from pathlib import Path
from typing import Dict, List, Any, Optional
from collections import OrderedDict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement


# Delagrave corporate blue color (extracted from template theme or standard corporate blue)
DELAGRAVE_BLUE = RGBColor(0, 85, 164)  # Corporate blue

# Fixed family ordering for document structure
FAMILY_ORDER = [
    'paillasse',
    'sorbonne',
    'revetement',
    'meubles',
    'tables-en',
    'equipement',
    'elec-sorb',
    'complements',
    'armoire-securite'
]

# Family display names (French with proper accents)
FAMILY_DISPLAY_NAMES = {
    'paillasse': 'Paillasses',
    'sorbonne': 'Sorbonnes',
    'revetement': 'Revêtements',
    'meubles': 'Meubles',
    'tables-en': 'Tables EN',
    'equipement': 'Équipement',
    'elec-sorb': 'Électricité Sorbonne',
    'complements': 'Compléments',
    'armoire-securite': 'Armoires de Sécurité'
}


def add_cover_page(prs: Presentation, devis_info: Dict[str, str], logo_path: Optional[Path] = None) -> None:
    """
    Add cover page with Delagrave branding and devis information.

    Cover page includes:
    - Blue bandeau at top (~15% height) with logo or "DELAGRAVE" text
    - Blue bandeau at bottom (~8% height)
    - Project title centered
    - Devis info (numero, date, client) below title

    Args:
        prs: Presentation object to add cover page to
        devis_info: Dictionary with 'numero_devis', 'date', 'client', 'titre_affaire'
        logo_path: Optional path to logo image (if None, uses text placeholder)
    """
    # Use blank layout (layout 0 or first available)
    blank_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)

    # Get slide dimensions (A4 portrait: 7561263 x 10693400 EMU)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Top bandeau (15% of slide height)
    top_bandeau_height = int(slide_height * 0.15)
    top_bandeau = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,  # left
        0,  # top
        slide_width,
        top_bandeau_height
    )
    top_bandeau.fill.solid()
    top_bandeau.fill.fore_color.rgb = DELAGRAVE_BLUE
    top_bandeau.line.fill.background()

    # Bottom bandeau (8% of slide height)
    bottom_bandeau_height = int(slide_height * 0.08)
    bottom_bandeau = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,  # left
        slide_height - bottom_bandeau_height,  # top
        slide_width,
        bottom_bandeau_height
    )
    bottom_bandeau.fill.solid()
    bottom_bandeau.fill.fore_color.rgb = DELAGRAVE_BLUE
    bottom_bandeau.line.fill.background()

    # Logo below bandeau on white background (logo is light blue, invisible on blue)
    if logo_path and logo_path.exists():
        logo_margin_left = int(slide_width * 0.03)
        logo_top = top_bandeau_height + int(slide_height * 0.02)
        logo_max_height = int(slide_height * 0.08)

        pic = slide.shapes.add_picture(
            str(logo_path),
            logo_margin_left,
            logo_top,
            height=logo_max_height
        )
    else:
        # Fallback: "DELAGRAVE" text in bandeau
        text_box = slide.shapes.add_textbox(
            int(slide_width * 0.05),
            int(top_bandeau_height * 0.25),
            int(slide_width * 0.5),
            int(top_bandeau_height * 0.5)
        )
        text_frame = text_box.text_frame
        text_frame.text = "DELAGRAVE"

        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # Project title (centered, below top bandeau)
    title_text = devis_info.get('titre_affaire', devis_info.get('client', 'Dossier de Fiches Techniques'))

    title_top = int(slide_height * 0.35)
    title_height = int(slide_height * 0.15)
    title_box = slide.shapes.add_textbox(
        int(slide_width * 0.1),
        title_top,
        int(slide_width * 0.8),
        title_height
    )

    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = title_text

    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_paragraph.font.size = Pt(24)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 0, 0)

    # Devis info (numero, date, client) below title
    info_lines = []
    if devis_info.get('numero_devis'):
        info_lines.append(f"Devis N° {devis_info['numero_devis']}")
    if devis_info.get('date'):
        info_lines.append(f"Date: {devis_info['date']}")
    if devis_info.get('client'):
        info_lines.append(f"Client: {devis_info['client']}")

    if info_lines:
        info_top = title_top + title_height + int(slide_height * 0.05)
        info_height = int(slide_height * 0.15)
        info_box = slide.shapes.add_textbox(
            int(slide_width * 0.1),
            info_top,
            int(slide_width * 0.8),
            info_height
        )

        info_frame = info_box.text_frame
        info_frame.word_wrap = True

        # Add each line as a separate paragraph
        for i, line in enumerate(info_lines):
            if i == 0:
                p = info_frame.paragraphs[0]
            else:
                p = info_frame.add_paragraph()

            p.text = line
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(100, 100, 100)


def add_chapter_separator(prs: Presentation, family_name: str) -> None:
    """
    Add chapter separator page for a product family.

    Separator includes:
    - Blue bandeau at top (~12% height) with family name in white
    - Rest of page is blank/white

    Args:
        prs: Presentation object to add separator to
        family_name: Internal family name (e.g., 'paillasse')
    """
    # Use blank layout
    blank_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Top bandeau (12% of slide height)
    bandeau_height = int(slide_height * 0.12)
    bandeau = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,  # left
        0,  # top
        slide_width,
        bandeau_height
    )
    bandeau.fill.solid()
    bandeau.fill.fore_color.rgb = DELAGRAVE_BLUE
    bandeau.line.fill.background()

    # Family name in white bold centered in bandeau
    family_display = FAMILY_DISPLAY_NAMES.get(family_name, family_name.title())

    text_box = slide.shapes.add_textbox(
        0,
        int(bandeau_height * 0.2),
        slide_width,
        int(bandeau_height * 0.6)
    )
    text_frame = text_box.text_frame
    text_frame.text = family_display

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(32)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)


def add_toc_page(prs: Presentation, toc_entries: List[Dict[str, Any]]) -> None:
    """
    Add table of contents page with families and product codes.

    TOC structure:
    - Title "Sommaire" at top
    - For each family:
      - Family name in bold (DELAGRAVE_BLUE)
      - Product codes with titles and page numbers

    Args:
        prs: Presentation object to add TOC to
        toc_entries: List of dicts with structure:
            [{'family': str, 'family_display': str,
              'products': [{'code': str, 'titre': str, 'page_number': int}]}]
    """
    # Use blank layout
    blank_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Title "Sommaire"
    title_height = int(slide_height * 0.08)
    title_box = slide.shapes.add_textbox(
        int(slide_width * 0.1),
        int(slide_height * 0.05),
        int(slide_width * 0.8),
        title_height
    )

    title_frame = title_box.text_frame
    title_frame.text = "Sommaire"

    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_paragraph.font.size = Pt(24)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = DELAGRAVE_BLUE

    # Content area
    content_top = int(slide_height * 0.15)
    content_height = int(slide_height * 0.75)
    content_width = int(slide_width * 0.8)
    content_box = slide.shapes.add_textbox(
        int(slide_width * 0.1),
        content_top,
        content_width,
        content_height
    )

    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Right tab stop position for page number alignment
    tab_pos_emu = content_width

    # Build TOC content
    first_paragraph = True
    for entry in toc_entries:
        family_display = entry['family_display']
        products = entry['products']

        if not products:
            continue

        # Family header (bold, blue)
        if first_paragraph:
            p = content_frame.paragraphs[0]
            first_paragraph = False
        else:
            p = content_frame.add_paragraph()

        p.text = family_display
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = DELAGRAVE_BLUE
        p.space_after = Pt(6)

        # Products (code + truncated title, page number right-aligned)
        for product in products:
            p = content_frame.add_paragraph()

            # Add right tab stop for page number alignment
            pPr = p._p.get_or_add_pPr()
            tabLst = OxmlElement('a:tabLst')
            tab = OxmlElement('a:tab')
            tab.set('pos', str(tab_pos_emu))
            tab.set('algn', 'r')
            tabLst.append(tab)
            pPr.append(tabLst)

            # Truncate title if too long
            titre = product['titre']
            if len(titre) > 50:
                titre = titre[:47] + '...'

            code = product['code']
            page_num = product['page_number']

            # Product text + tab
            run1 = p.add_run()
            run1.text = f"  {code} - {titre}\t"
            run1.font.size = Pt(11)

            # Page number (right-aligned via tab stop)
            run2 = p.add_run()
            run2.text = str(page_num)
            run2.font.size = Pt(11)
            run2.font.bold = True

            p.space_after = Pt(3)


def add_page_numbers(prs: Presentation) -> None:
    """
    Add page numbers to all slides in the presentation.

    Adds a small text box at bottom-right of each slide with "P. X" format.
    """
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    box_width = Cm(2)
    box_height = Cm(0.7)
    left = slide_width - box_width - Cm(0.8)
    top = slide_height - box_height - Cm(0.3)

    for i, slide in enumerate(prs.slides):
        page_num = i + 1
        txBox = slide.shapes.add_textbox(left, top, box_width, box_height)
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"P. {page_num}"
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(100, 100, 100)


def assemble_document(
    product_groups: OrderedDict,
    prs: Presentation,
    devis_info: Dict[str, str],
    references_dir: Path,
    project_root: Path,
    logo_path: Optional[Path] = None
) -> Dict[str, Any]:
    """
    Assemble complete document with cover, TOC, chapter separators, and product slides.

    This is the main orchestrator that builds the entire document structure:
    1. First pass: calculate page numbers for all content
    2. Add cover page (slide 1)
    3. Add TOC with calculated page numbers (slide 2)
    4. For each family in FAMILY_ORDER:
       - Add chapter separator
       - Add product slides for that family

    Args:
        product_groups: OrderedDict mapping family names to lists of product dicts
            Each product dict should have: 'code', 'titre', 'famille', and full product data
        prs: Presentation object to build document in
        devis_info: Dictionary with 'numero_devis', 'date', 'client'
        references_dir: Path to Delagrave/references/
        project_root: Root directory for resolving image paths
        logo_path: Optional path to logo image

    Returns:
        Dictionary with assembly statistics:
        {
            'total_pages': int,
            'families_included': list of family names,
            'product_count': int,
            'toc_entries': list of TOC entry dicts
        }
    """
    # Use modern template for all slide building
    from gendoc.generators.modern_template import (
        build_cover, build_toc, build_separator,
        build_product_slide, add_page_numbers as modern_page_numbers
    )

    # Step 1: First pass - calculate page numbers and build TOC entries
    page_counter = 2  # Cover=1, TOC=2, content starts at 3
    toc_entries = []
    families_included = []

    for family in FAMILY_ORDER:
        if family not in product_groups or not product_groups[family]:
            continue

        families_included.append(family)

        # Chapter separator will be one page
        page_counter += 1

        # Track products for TOC
        family_products = []
        products = product_groups[family]

        for product in products:
            page_counter += 1  # Each product gets a slide

            # Record for TOC
            family_products.append({
                'code': product.get('code', ''),
                'titre': product.get('titre', ''),
                'page_number': page_counter
            })

        # Add family to TOC entries
        if family_products:
            toc_entries.append({
                'family': family,
                'family_display': FAMILY_DISPLAY_NAMES.get(family, family.title()),
                'products': family_products
            })

    total_pages = page_counter
    product_count = sum(len(products) for products in product_groups.values())

    # Step 2: Build slides using modern template
    build_cover(prs, devis_info, logo_path)
    build_toc(prs, toc_entries)

    # Step 3: Add content slides
    all_warnings = []
    for family in FAMILY_ORDER:
        if family not in product_groups or not product_groups[family]:
            continue

        # Add chapter separator
        family_display = FAMILY_DISPLAY_NAMES.get(family, family.title())
        build_separator(prs, family, family_display)

        # Add product slides
        for product in product_groups[family]:
            slide_warnings = build_product_slide(prs, product, family, project_root, logo_path)
            if slide_warnings:
                for w in slide_warnings:
                    all_warnings.append({"code": product.get("code", "?"), "message": w})

    # Page numbers added later by generate_presentation (after revetement slides)

    # Return assembly statistics
    return {
        'total_pages': total_pages,
        'families_included': families_included,
        'product_count': product_count,
        'toc_entries': toc_entries,
        'warnings': all_warnings
    }
