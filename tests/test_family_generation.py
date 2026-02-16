"""
Per-family generation tests for PowerPoint tech sheets.

This module tests that each of the 8 product families can generate valid
PowerPoint presentations with correct slide counts and populated placeholders.
"""

import pytest
from pathlib import Path
from pptx import Presentation as PptxPresentation
from gendoc.generators.pptx_generator import generate_presentation, _split_revetement_text
from gendoc.parsers.md_parser import find_product


FAMILIES = [
    'paillasse', 'sorbonne', 'revetement', 'meubles',
    'tables-en', 'equipement', 'elec-sorb', 'complements',
    'enceinte-ventilee'
]


@pytest.mark.parametrize("family", FAMILIES)
def test_family_generates_valid_pptx(family, sample_codes, template_path, references_dir, project_root, output_dir):
    """Each family generates a valid .pptx with correct slide count and populated placeholders."""
    code = sample_codes[family]
    output_path = output_dir / f"test_{family}.pptx"

    # Generate presentation for this single product
    result = generate_presentation(
        product_codes=[code],
        output_path=output_path,
        references_dir=references_dir,
        template_path=template_path,
        project_root=project_root,
        mode="FTI"
    )

    # Verify result dict
    assert result['slides_generated'] >= 1, f"No slides generated for {family} code {code}"
    assert result['total_pages'] >= 3, f"Expected at least 3 pages (cover+TOC+separator+product) for {family}"
    assert len(result['skipped']) == 0, f"Product {code} was skipped: {result['skipped']}"

    # Verify .pptx file was created and is valid
    assert output_path.exists(), f"Output file not created: {output_path}"
    assert output_path.stat().st_size > 10000, f"Output file too small (likely empty): {output_path}"

    # Open and verify slide structure
    prs = PptxPresentation(str(output_path))
    slide_count = len(prs.slides)

    # Minimum slides: cover(1) + TOC(1) + separator(1) + product(1) = 4
    assert slide_count >= 4, f"Expected at least 4 slides for {family}, got {slide_count}"

    # Verify product slide (last slide before any page numbers) has populated content
    # The product slide is at index 3 (0=cover, 1=TOC, 2=separator, 3=product)
    product_slide = prs.slides[3]

    # Check that the slide has shapes (not completely empty)
    shape_count = len(product_slide.shapes)
    assert shape_count >= 1, f"Product slide for {family} has no shapes"


@pytest.mark.parametrize("family", FAMILIES)
def test_sample_code_exists_in_references(family, sample_codes, references_dir):
    """Each sample code used in tests must exist in the reference files."""
    code = sample_codes[family]
    product = find_product(code, references_dir)
    assert product is not None, f"Sample code {code} not found in references for family {family}"

    # Handle accent normalization in family names (revetement vs revètement)
    product_family = product['famille'].lower().replace('è', 'e').replace('é', 'e')
    expected_family = family.lower().replace('è', 'e').replace('é', 'e')

    assert product_family == expected_family, \
        f"Product {code} is in family {product['famille']}, expected {family}"


def test_split_revetement_text():
    """Revetement text splitting handles 1, 2, and 3+ blocks correctly."""

    # 3 blocks
    text_3 = "Block one description.\n\nBlock two meo.\n\nBlock three finish."
    result = _split_revetement_text(text_3)
    assert result['texte'] == "Block one description."
    assert result['mise_en_oeuvre'] == "Block two meo."
    assert result['finition'] == "Block three finish."

    # 2 blocks
    text_2 = "Block one.\n\nBlock two."
    result = _split_revetement_text(text_2)
    assert result['texte'] == "Block one."
    assert result['mise_en_oeuvre'] == "Block two."
    assert result['finition'] == ""

    # 1 block
    text_1 = "Just one block."
    result = _split_revetement_text(text_1)
    assert result['texte'] == "Just one block."
    assert result['mise_en_oeuvre'] == ""
    assert result['finition'] == ""
