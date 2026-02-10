"""
End-to-end pipeline tests for Delagrave quote analysis and presentation generation.

These tests validate the complete user workflow:
1. Analyze a quote PDF (extract header and product references)
2. Generate a PowerPoint presentation from the extracted codes
3. Verify the generated PowerPoint is valid and contains expected slides
"""

import pytest
from pathlib import Path
from pptx import Presentation as PptxPresentation

from gendoc.parsers.devis_analyzer import analyze_devis
from gendoc.generators.pptx_generator import generate_presentation


class TestE2EPipeline:
    """End-to-end pipeline: analyze devis PDF -> generate PowerPoint."""

    def test_analyze_devis_extracts_references(self, project_root, references_dir):
        """Step 1: analyze_devis extracts header and product references from PDF."""
        pdf_path = project_root / "Delagrave" / "Devis - Modeles" / "Devis Test.pdf"
        if not pdf_path.exists():
            pytest.skip(f"Test devis PDF not found: {pdf_path}")

        result = analyze_devis(pdf_path, references_dir)

        # Must have header info
        assert 'header' in result
        assert 'references' in result

        # Must find at least some references
        assert len(result['references']) > 0, "No references extracted from test devis"

        # References must have code and famille
        for ref in result['references']:
            assert 'code' in ref, f"Reference missing 'code': {ref}"
            assert 'famille' in ref, f"Reference missing 'famille': {ref}"

    def test_analyze_devis_classifies_families(self, project_root, references_dir):
        """Step 2: classified references span multiple families."""
        pdf_path = project_root / "Delagrave" / "Devis - Modeles" / "Devis Test.pdf"
        if not pdf_path.exists():
            pytest.skip(f"Test devis PDF not found: {pdf_path}")

        result = analyze_devis(pdf_path, references_dir)

        families = set(ref['famille'] for ref in result['references'])
        assert len(families) >= 2, f"Expected references from at least 2 families, got: {families}"

    def test_full_pipeline_analyze_then_generate(self, project_root, references_dir, template_path, output_dir):
        """Step 3: Full E2E - analyze devis then generate PowerPoint from extracted codes."""
        pdf_path = project_root / "Delagrave" / "Devis - Modeles" / "Devis Test.pdf"
        if not pdf_path.exists():
            pytest.skip(f"Test devis PDF not found: {pdf_path}")

        # Phase 1: Analyze
        analysis = analyze_devis(pdf_path, references_dir)
        product_codes = [ref['code'] for ref in analysis['references']]
        revetement_codes = [rev['code'] for rev in analysis.get('revetements', [])]

        assert len(product_codes) > 0, "No product codes from analysis"

        # Phase 2: Generate
        output_path = output_dir / "test_e2e_pipeline.pptx"
        result = generate_presentation(
            product_codes=product_codes,
            output_path=output_path,
            references_dir=references_dir,
            template_path=template_path,
            project_root=project_root,
            mode="FTI",
            revetement_codes=revetement_codes,
            devis_info=analysis.get('header', {})
        )

        # Verify generation succeeded
        assert result['slides_generated'] > 0, f"No slides generated. Skipped: {result['skipped']}"
        assert output_path.exists(), f"Output file not created: {output_path}"

        # Verify the file is a valid PowerPoint
        prs = PptxPresentation(str(output_path))
        slide_count = len(prs.slides)

        # Minimum: cover + TOC + at least one family (separator + product) = 4
        assert slide_count >= 4, f"Expected at least 4 slides, got {slide_count}"

        # Verify cover slide exists (first slide)
        cover_slide = prs.slides[0]
        assert len(cover_slide.shapes) >= 2, "Cover slide seems empty"

    def test_pipeline_with_sp_devis(self, project_root, references_dir, template_path, output_dir):
        """E2E with SP (special articles) devis if available."""
        pdf_path = project_root / "Delagrave" / "Devis - Modeles" / "Devis avec SP.pdf"
        if not pdf_path.exists():
            pytest.skip(f"SP devis PDF not found: {pdf_path}")

        analysis = analyze_devis(pdf_path, references_dir)

        # Should detect some special articles
        assert 'speciaux' in analysis, "Analysis missing 'speciaux' key"

        # The SP devis should have at least regular references
        product_codes = [ref['code'] for ref in analysis['references']]
        if len(product_codes) == 0:
            pytest.skip("No regular product codes found in SP devis")

        output_path = output_dir / "test_e2e_sp_pipeline.pptx"
        result = generate_presentation(
            product_codes=product_codes,
            output_path=output_path,
            references_dir=references_dir,
            template_path=template_path,
            project_root=project_root,
            mode="FTI",
            devis_info=analysis.get('header', {})
        )

        assert result['slides_generated'] > 0, "No slides generated from SP devis"
        assert output_path.exists()
